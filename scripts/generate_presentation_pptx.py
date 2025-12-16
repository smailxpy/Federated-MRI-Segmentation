"""
Generate a professional PPTX presentation for the project report.

- Uses metrics from: results/final_experiment/final_report.json
- Uses figures from: results/figures/*.png

Output:
  - Federated_Continual_Learning_for_MRI_Brain_Tumor_Segmentation_Team.pptx (repo root)
"""

from __future__ import annotations

import argparse
import json
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
REPORT_JSON = REPO_ROOT / "results" / "final_experiment" / "final_report.json"
FIG_DIR = REPO_ROOT / "results" / "figures"
OUT_PPTX = REPO_ROOT / "Federated_Continual_Learning_for_MRI_Brain_Tumor_Segmentation_Team.pptx"


def _load_json(path: Path) -> Dict[str, Any]:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def _pct(x: float, digits: int = 2) -> str:
    return f"{x * 100:.{digits}f}%"


def _mm(x: float, digits: int = 2) -> str:
    return f"{x:.{digits}f} mm"


def _set_slide_bg_white(slide) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def _add_title(slide, title: str, x=Inches(0.6), y=Inches(0.3), w=Inches(12.1), h=Inches(0.7)) -> None:
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Times New Roman"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT


def _add_subtitle(slide, subtitle: str, x=Inches(0.6), y=Inches(1.05), w=Inches(12.1), h=Inches(0.6)) -> None:
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Times New Roman"
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(60, 60, 60)
    p.alignment = PP_ALIGN.LEFT


def _add_bullets(slide, bullets, x, y, w, h, font_size=Pt(18)) -> None:
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(b)
        p.level = 0
        p.font.name = "Times New Roman"
        p.font.size = font_size
        p.font.color.rgb = RGBColor(0, 0, 0)


def _add_table(slide, rows: int, cols: int, x, y, w, h, header: Tuple[str, ...], data_rows) -> None:
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = table_shape.table

    # Header style
    for j in range(cols):
        cell = tbl.cell(0, j)
        cell.text = header[j]
        for p in cell.text_frame.paragraphs:
            p.font.name = "Times New Roman"
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(0, 0, 0)

    # Body
    for i, r in enumerate(data_rows, start=1):
        for j in range(cols):
            cell = tbl.cell(i, j)
            cell.text = str(r[j])
            for p in cell.text_frame.paragraphs:
                p.font.name = "Times New Roman"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0)


def _add_picture_if_exists(slide, path: Path, x, y, w=None, h=None) -> bool:
    if not path.exists():
        return False
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    elif w is not None:
        slide.shapes.add_picture(str(path), x, y, width=w)
    elif h is not None:
        slide.shapes.add_picture(str(path), x, y, height=h)
    else:
        slide.shapes.add_picture(str(path), x, y)
    return True


def build_presentation(report: Dict[str, Any]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    team = report.get("team", "314IV")
    topic = report.get("topic", "Federated Continual Learning for MRI Segmentation")
    ts = report.get("timestamp")
    try:
        date_str = datetime.fromisoformat(ts).strftime("%b %d, %Y") if ts else "Dec 2025"
    except Exception:
        date_str = "Dec 2025"

    dataset = report["dataset"]
    training = report["training"]
    ttime = report["training_time"]
    res = report["results"]
    overall = res["overall_metrics"]
    per = res["per_class_metrics"]
    clm = res["continual_learning_metrics"]

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Federated Continual Learning for MRI Brain Tumor Segmentation")
    _add_subtitle(slide, f"Team {team} | Topic: {topic} | {date_str}")
    _add_bullets(
        slide,
        [
            "Model: SegResNet + drift-aware adapters",
            "Federation: 4 hospitals (FedAvg)",
            f"Dataset: BraTS2021 subset (n={dataset['total_patients']})",
        ],
        x=Inches(0.8),
        y=Inches(2.0),
        w=Inches(12.0),
        h=Inches(4.8),
        font_size=Pt(20),
    )

    # Slide 2: Problem & Motivation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Problem & Motivation")
    _add_bullets(
        slide,
        [
            "Goal: segment brain tumor regions from 3D multi-modal MRI.",
            "Constraint 1: hospitals cannot share raw patient data (privacy).",
            "Constraint 2: distribution shift across hospitals and over time.",
            "Need: federated training + continual adaptation with limited forgetting.",
        ],
        x=Inches(0.8),
        y=Inches(1.6),
        w=Inches(12.0),
        h=Inches(5.6),
        font_size=Pt(20),
    )

    # Slide 3: Dataset & Federated Split
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Dataset & Federated Split")
    _add_table(
        slide,
        rows=5,
        cols=2,
        x=Inches(0.8),
        y=Inches(1.6),
        w=Inches(6.2),
        h=Inches(3.0),
        header=("Property", "Value"),
        data_rows=[
            ("Dataset", str(dataset["name"])),
            ("Patients", f"{dataset['total_patients']} (480 train / 60 val / 60 test)"),
            ("Modalities", ", ".join(dataset["modalities"])),
            ("Input size", "224×224×144 (voxel spacing 1mm)"),
        ],
    )
    _add_bullets(
        slide,
        [
            "Federated setup: 4 simulated hospitals",
            "120 patients per hospital (training pool)",
            "Heterogeneity: scanner/protocol variations (simulated by splits)",
        ],
        x=Inches(7.4),
        y=Inches(1.6),
        w=Inches(5.6),
        h=Inches(3.0),
        font_size=Pt(18),
    )

    # Slide 4: Method Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Method Overview")
    _add_bullets(
        slide,
        [
            "Backbone: 3D SegResNet (residual encoder–decoder).",
            "Client-side adaptation: drift-aware adapters (lightweight modules).",
            "Federated aggregation: FedAvg across 4 clients.",
            "Objective: improve robustness under domain shift and reduce forgetting.",
        ],
        x=Inches(0.8),
        y=Inches(1.6),
        w=Inches(12.0),
        h=Inches(5.6),
        font_size=Pt(20),
    )

    # Slide 5: Training Setup (RTX 5070)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Training Setup (Single RTX 5070)")
    gpu = report["hardware"]["gpu"]
    _add_table(
        slide,
        rows=6,
        cols=2,
        x=Inches(0.8),
        y=Inches(1.6),
        w=Inches(6.2),
        h=Inches(3.6),
        header=("Setting", "Value"),
        data_rows=[
            ("GPU", f"{gpu['model']} ({gpu['vram']})"),
            ("Batch size", str(training["batch_size"])),
            ("Rounds", f"{training['num_rounds']} (early stop @ {training.get('rounds_completed', 185)})"),
            ("Local epochs/round", str(training["local_epochs"])),
            ("LR / Optimizer", f"{training['learning_rate']} / {training['optimizer']}"),
        ],
    )
    _add_bullets(
        slide,
        [
            f"Total training time: ~{ttime['total_hours']} hours",
            f"Per round: {ttime['time_per_round']}",
            f"AMP: {str(training.get('amp_enabled', True)).lower()}",
            "Checkpointing enabled (best + latest).",
        ],
        x=Inches(7.4),
        y=Inches(1.6),
        w=Inches(5.6),
        h=Inches(3.6),
        font_size=Pt(18),
    )

    # Slide 6: Results Summary Table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Results (Test Set, n=60)")
    _add_table(
        slide,
        rows=5,
        cols=5,
        x=Inches(0.6),
        y=Inches(1.6),
        w=Inches(12.2),
        h=Inches(2.4),
        header=("Metric", "TC", "WT", "ET", "Mean"),
        data_rows=[
            ("Dice", _pct(per["tumor_core_TC"]["dice"]), _pct(per["whole_tumor_WT"]["dice"]), _pct(per["enhancing_tumor_ET"]["dice"]), _pct(overall["dice_mean"])),
            ("IoU", _pct(per["tumor_core_TC"]["iou"]), _pct(per["whole_tumor_WT"]["iou"]), _pct(per["enhancing_tumor_ET"]["iou"]), _pct(overall["iou_mean"])),
            ("HD95", _mm(per["tumor_core_TC"]["hd95_mm"]), _mm(per["whole_tumor_WT"]["hd95_mm"]), _mm(per["enhancing_tumor_ET"]["hd95_mm"]), _mm(overall["hd95_mean_mm"])),
            ("ASSD", _mm(per["tumor_core_TC"]["assd_mm"]), _mm(per["whole_tumor_WT"]["assd_mm"]), _mm(per["enhancing_tumor_ET"]["assd_mm"]), _mm(overall["assd_mean_mm"])),
        ],
    )
    _add_bullets(
        slide,
        [
            f"Mean Dice: {_pct(overall['dice_mean'])} (CI95: {_pct(res['statistical_analysis']['confidence_interval_95']['dice_lower'])}–{_pct(res['statistical_analysis']['confidence_interval_95']['dice_upper'])})",
            f"Avg forgetting rate: {_pct(clm['average_forgetting_rate'])}",
            f"Backward transfer: {_pct(clm['backward_transfer'])} | Forward transfer: {_pct(clm['forward_transfer'])}",
            "Metrics: Dice, IoU, HD95, ASSD, Sensitivity, Precision, Specificity.",
        ],
        x=Inches(0.8),
        y=Inches(4.2),
        w=Inches(12.0),
        h=Inches(2.8),
        font_size=Pt(18),
    )

    # Slide 7: Training curves (Fig1 + Fig2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Training Dynamics")
    _add_picture_if_exists(slide, FIG_DIR / "fig1_training_progression.png", Inches(0.6), Inches(1.4), w=Inches(6.3))
    _add_picture_if_exists(slide, FIG_DIR / "fig2_loss_curve.png", Inches(6.8), Inches(1.4), w=Inches(6.3))
    _add_subtitle(slide, "Dice progression and loss convergence (early stopping at ~round 185)", y=Inches(6.95), h=Inches(0.4))

    # Slide 8: Baseline comparison + per-class metrics (Fig3 + Fig4)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Baselines & Per-Class Performance")
    _add_picture_if_exists(slide, FIG_DIR / "fig3_method_comparison.png", Inches(0.6), Inches(1.4), w=Inches(6.3))
    _add_picture_if_exists(slide, FIG_DIR / "fig4_per_class_metrics.png", Inches(6.8), Inches(1.4), w=Inches(6.3))
    _add_bullets(
        slide,
        [
            "Gap vs centralized is expected in federated training (privacy preserved).",
            "WT usually easiest; ET hardest due to small size and imbalance.",
        ],
        x=Inches(0.9),
        y=Inches(6.35),
        w=Inches(12.0),
        h=Inches(1.0),
        font_size=Pt(16),
    )

    # Slide 9: Distribution + forgetting (Fig5 + Fig7)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Generalization & Continual Learning")
    _add_picture_if_exists(slide, FIG_DIR / "fig5_dice_distribution.png", Inches(0.6), Inches(1.4), w=Inches(6.3))
    _add_picture_if_exists(slide, FIG_DIR / "fig7_forgetting_analysis.png", Inches(6.8), Inches(1.4), w=Inches(6.3))
    _add_bullets(
        slide,
        [
            f"Dice distribution (n={dataset['test_samples']}) shows stable performance.",
            f"Avg forgetting: {_pct(clm['average_forgetting_rate'])} (low-to-moderate).",
        ],
        x=Inches(0.9),
        y=Inches(6.35),
        w=Inches(12.0),
        h=Inches(1.0),
        font_size=Pt(16),
    )

    # Slide 10: Qualitative + Conclusion (Fig8)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_title(slide, "Qualitative Results & Conclusion")
    _add_picture_if_exists(slide, FIG_DIR / "fig8_segmentation_examples.png", Inches(0.6), Inches(1.5), w=Inches(7.1))
    _add_bullets(
        slide,
        [
            f"Achieved {_pct(overall['dice_mean'])} mean Dice on BraTS2021 subset (target ≥ 70%).",
            f"Single-GPU training (RTX 5070, 12GB): ~{ttime['total_hours']} hours with early stopping.",
            f"Continual learning: avg forgetting {_pct(clm['average_forgetting_rate'])}, BWT {_pct(clm['backward_transfer'])}.",
            f"Boundary quality: mean HD95 {_mm(overall['hd95_mean_mm'])}, mean ASSD {_mm(overall['assd_mean_mm'])}.",
            "Future work: uncertainty estimation, stronger ET-focused losses, real multi-institution evaluation.",
        ],
        x=Inches(7.9),
        y=Inches(1.6),
        w=Inches(5.1),
        h=Inches(5.6),
        font_size=Pt(18),
    )

    return prs


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=str, default=str(OUT_PPTX), help="Output .pptx path")
    args = parser.parse_args()

    if not REPORT_JSON.exists():
        raise FileNotFoundError(f"Missing report json: {REPORT_JSON}")
    if not FIG_DIR.exists():
        raise FileNotFoundError(f"Missing figures dir: {FIG_DIR}")

    report = _load_json(REPORT_JSON)
    prs = build_presentation(report)
    out_path = Path(args.out)
    prs.save(str(out_path))
    print(f"[OK] Wrote PPTX: {out_path}")


if __name__ == "__main__":
    main()


